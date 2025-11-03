# shared/document_processor.py
"""
Universal Document Processor
Handles multi-format document text extraction
"""

import io
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from pathlib import Path
from typing import Union, BinaryIO
import logging

class UniversalDocumentProcessor:
    """Process multiple document formats and extract text"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def process_document_from_bytes(self, file_data: bytes, filename: str) -> str:
        """
        Process document from bytes based on file extension
        
        Args:
            file_data: Raw file bytes
            filename: Original filename for extension detection
            
        Returns:
            Extracted text content
        """
        suffix = Path(filename).suffix.lower()
        
        processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_word,
            '.doc': self._process_word,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.csv': self._process_csv,
            '.txt': self._process_text,
        }
        
        if suffix in processors:
            return processors[suffix](file_data)
        else:
            self.logger.warning(f"Unsupported file type: {suffix}")
            raise ValueError(f"Unsupported file type: {suffix}")
    
    def _process_pdf(self, file_data: bytes) -> str:
        """Extract text from PDF using PyMuPDF"""
        text_parts = []
        
        try:
            with fitz.open(stream=file_data, filetype="pdf") as doc:
                for page_num, page in enumerate(doc, 1):
                    page_text = page.get_text()
                    if page_text.strip():
                        text_parts.append(f"[Page {page_num}]\n{page_text}")
                    
                    # Extract tables if present
                    tables = page.find_tables()
                    if tables:
                        for table in tables:
                            df = table.to_pandas()
                            if not df.empty:
                                text_parts.append(f"\n{df.to_markdown(index=False)}\n")
        
        except Exception as e:
            self.logger.error(f"Error processing PDF: {str(e)}")
            raise
        
        return "\n\n".join(text_parts)
    
    def _process_word(self, file_data: bytes) -> str:
        """Extract text from Word document"""
        text_parts = []
        
        try:
            with io.BytesIO(file_data) as file_stream:
                doc = Document(file_stream)
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                
                # Extract tables as markdown
                for table in doc.tables:
                    table_data = [
                        [cell.text for cell in row.cells]
                        for row in table.rows
                    ]
                    if table_data:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        text_parts.append(f"\n{df.to_markdown(index=False)}\n")
        
        except Exception as e:
            self.logger.error(f"Error processing Word document: {str(e)}")
            raise
        
        return "\n\n".join(text_parts)
    
    def _process_excel(self, file_data: bytes) -> str:
        """Extract data from Excel spreadsheet"""
        text_parts = []
        
        try:
            with io.BytesIO(file_data) as file_stream:
                excel_file = pd.ExcelFile(file_stream)
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    if not df.empty:
                        text_parts.append(f"## Sheet: {sheet_name}")
                        text_parts.append(df.to_markdown(index=False))
                        text_parts.append(f"\n**Summary:** {len(df)} rows, {len(df.columns)} columns\n")
        
        except Exception as e:
            self.logger.error(f"Error processing Excel file: {str(e)}")
            raise
        
        return "\n\n".join(text_parts)
    
    def _process_powerpoint(self, file_data: bytes) -> str:
        """Extract text from PowerPoint presentation"""
        text_parts = []
        
        try:
            with io.BytesIO(file_data) as file_stream:
                prs = Presentation(file_stream)
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = [f"[Slide {slide_num}]"]
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                        
                        # Extract tables
                        if shape.has_table:
                            table_data = [
                                [cell.text for cell in row.cells]
                                for row in shape.table.rows
                            ]
                            if table_data:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                slide_text.append(df.to_markdown(index=False))
                    
                    text_parts.append("\n".join(slide_text))
        
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint: {str(e)}")
            raise
        
        return "\n\n".join(text_parts)
    
    def _process_csv(self, file_data: bytes) -> str:
        """Extract data from CSV file"""
        try:
            with io.BytesIO(file_data) as file_stream:
                df = pd.read_csv(file_stream)
                return df.to_markdown(index=False)
        except Exception as e:
            self.logger.error(f"Error processing CSV: {str(e)}")
            raise
    
    def _process_text(self, file_data: bytes) -> str:
        """Extract text from plain text file"""
        try:
            return file_data.decode('utf-8')
        except UnicodeDecodeError:
            # Try other encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    return file_data.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            self.logger.error("Could not decode text file with any encoding")
            raise ValueError("Unable to decode text file")
